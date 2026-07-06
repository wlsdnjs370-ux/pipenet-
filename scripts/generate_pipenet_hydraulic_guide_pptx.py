from __future__ import annotations

import os
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DESKTOP = Path(os.environ.get("USERPROFILE", r"C:\Users\admin")) / "Desktop"
OUT = DESKTOP / "PIPENET 수리계산 검증 프로그램 가이드_v3_편집가능.pptx"
OUT_ASCII = DESKTOP / "PIPENET_hydraulic_calculation_guide_v3_editable.pptx"
SCREENSHOT_DIR = ROOT / "docs" / "generated_reports" / "server_screenshots"


W, H = 11.69, 8.27
SIDEBAR_W = 3.05
LEFT = SIDEBAR_W + 0.42
RIGHT = W - 0.42
CONTENT_W = RIGHT - LEFT
FONT = "함초롬바탕"

NAVY = RGBColor(18, 35, 58)
BLUE = RGBColor(37, 99, 235)
SKY = RGBColor(226, 238, 255)
GREEN = RGBColor(22, 128, 61)
GREEN_BG = RGBColor(220, 252, 231)
ORANGE = RGBColor(194, 65, 12)
ORANGE_BG = RGBColor(255, 237, 213)
RED = RGBColor(220, 38, 38)
RED_BG = RGBColor(254, 226, 226)
GRAY = RGBColor(241, 245, 249)
TEXT = RGBColor(15, 23, 42)
MUTED = RGBColor(100, 116, 139)
WHITE = RGBColor(255, 255, 255)


def add_text(slide, x, y, w, h, text, size=12, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_rect(slide, x, y, w, h, fill, line=None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    shape.line.width = Pt(0.75)
    return shape


def add_card(slide, x, y, w, h, title, body, accent=BLUE, fill=WHITE):
    add_rect(slide, x, y, w, h, fill, RGBColor(203, 213, 225))
    add_rect(slide, x, y, 0.08, h, accent, accent)
    add_text(slide, x + 0.16, y + 0.12, w - 0.28, 0.28, title, 11, accent, True)
    add_text(slide, x + 0.16, y + 0.45, w - 0.28, h - 0.55, body, 8.3, TEXT)


def add_sidebar(slide, no, title, keywords, note_title, notes):
    add_rect(slide, 0, 0, SIDEBAR_W, H, NAVY, NAVY)
    add_text(slide, 0.32, 0.38, 2.3, 0.3, "PIPENET GUIDE", 14, WHITE, True)
    add_text(slide, 0.32, 0.85, 1.4, 0.45, f"{no:02d}", 28, RGBColor(166, 191, 255), True)
    add_text(slide, 0.32, 1.45, 2.25, 0.7, title, 12, WHITE, True)
    add_text(slide, 0.32, 2.55, 2.1, 0.25, "KEY POINTS", 9.5, RGBColor(166, 191, 255), True)
    y = 2.95
    for kw in keywords:
        add_rect(slide, 0.34, y + 0.06, 0.06, 0.06, RGBColor(166, 191, 255), RGBColor(166, 191, 255))
        add_text(slide, 0.48, y, 2.15, 0.25, kw, 8, WHITE)
        y += 0.35
    add_rect(slide, 0.32, 6.05, 2.28, 1.55, RGBColor(31, 41, 64), RGBColor(63, 73, 94))
    add_text(slide, 0.48, 6.22, 2.0, 0.25, note_title, 9, RGBColor(166, 191, 255), True)
    add_text(slide, 0.48, 6.58, 1.95, 0.78, "\n".join(notes), 7.5, WHITE)


def add_header(slide, no, title, subtitle):
    add_text(slide, LEFT, 0.38, CONTENT_W, 0.36, f"{no:02d}. {title}", 19, TEXT, True)
    add_text(slide, LEFT, 0.82, CONTENT_W, 0.28, subtitle, 9, MUTED)


def add_table(slide, x, y, widths, row_h, headers, rows, title=None):
    if title:
        add_text(slide, x, y - 0.32, sum(widths), 0.24, title, 10.5, TEXT, True)
    cols = len(headers)
    rows_count = len(rows) + 1
    table_shape = slide.shapes.add_table(rows_count, cols, Inches(x), Inches(y), Inches(sum(widths)), Inches(row_h * rows_count))
    table = table_shape.table
    for i, width in enumerate(widths):
        table.columns[i].width = Inches(width)
    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.runs[0].font.name = FONT
        p.runs[0].font.size = Pt(7.5)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = WHITE
    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 else GRAY
            p = cell.text_frame.paragraphs[0]
            p.runs[0].font.name = FONT
            p.runs[0].font.size = Pt(6.6)
            p.runs[0].font.color.rgb = TEXT


def add_flow(slide, x, y, labels, box_w=1.2, box_h=0.75, gap=0.22):
    cx = x
    for i, (head, body) in enumerate(labels):
        add_rect(slide, cx, y, box_w, box_h, SKY, BLUE)
        add_text(slide, cx + 0.08, y + 0.08, box_w - 0.16, 0.18, head, 8, BLUE, True, PP_ALIGN.CENTER)
        add_text(slide, cx + 0.08, y + 0.33, box_w - 0.16, 0.28, body, 6.4, TEXT, False, PP_ALIGN.CENTER)
        if i < len(labels) - 1:
            add_text(slide, cx + box_w + 0.03, y + 0.25, gap - 0.04, 0.2, "→", 14, BLUE, True, PP_ALIGN.CENTER)
        cx += box_w + gap


def add_image(slide, x, y, w, h, name, caption):
    path = SCREENSHOT_DIR / name
    add_rect(slide, x, y, w, h, WHITE, RGBColor(203, 213, 225))
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(x + 0.06), Inches(y + 0.06), width=Inches(w - 0.12), height=Inches(h - 0.34))
    add_text(slide, x + 0.06, y + h - 0.25, w - 0.12, 0.18, caption, 6.6, MUTED, False, PP_ALIGN.CENTER)


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)

    # 1 Cover
    s = blank_slide(prs)
    add_sidebar(s, 1, "표지", ["검증 로직", "결과 해석", "설계 최적화", "CAD 대조"], "DOCUMENT SCOPE", ["수리계산 결과 검증 기준", "데이터 테이블 해석", "SDF 배관망 시각화", "설계 최적화 리포트"])
    add_text(s, LEFT, 0.6, CONTENT_W, 0.45, "PIPENET 수리계산 검증 프로그램", 24, TEXT, True)
    add_text(s, LEFT, 1.12, CONTENT_W, 0.35, "검증 로직 · 결과 해석 · 설계 최적화 가이드", 15, BLUE, True)
    add_text(s, LEFT, 1.85, CONTENT_W, 0.8, "결과서 DOCX/PDF와 SDF 배관망을 기반으로 수리계산 결과의 법적 기준 충족 여부, Hazen-Williams 재계산, 배관 토폴로지 유속 판정, 공학/경제 최적화 후보를 확인하는 실무용 안내서입니다.", 10, TEXT)
    add_card(s, LEFT, 3.2, 1.65, 1.15, "검증", "헤드, 배관 유속, C-Factor, 특수설비, HW 마찰손실 재계산", BLUE, SKY)
    add_card(s, LEFT + 1.9, 3.2, 1.65, 1.15, "해석", "결과 데이터 테이블, 상세 카드, 색상 범례, 판정 사유 확인", GREEN, GREEN_BG)
    add_card(s, LEFT + 3.8, 3.2, 1.65, 1.15, "최적화", "공학적 안정성 후보와 시공사 경제성 후보를 분리 제시", ORANGE, ORANGE_BG)
    add_flow(s, LEFT, 5.35, [("입력", "DOCX/PDF\nSDF"), ("파싱", "표/그래프\n구조화"), ("검증", "PASS/FAIL\nREVIEW"), ("최적화", "공학/경제\n후보"), ("출력", "테이블\n맵")], 1.02, 0.78, 0.16)

    # 2 Data
    s = blank_slide(prs)
    add_sidebar(s, 2, "입력 데이터와 사용 흐름", ["결과서 DOCX/PDF", "SDF 배관망", "선택 입력", "결과 테이블"], "DATA FLOW", ["입력값은 원자료를 보존한 상태로", "검증용 구조 데이터로 변환됩니다."])
    add_header(s, 2, "입력 데이터와 사용 흐름", "결과서와 SDF 파일이 서버 내부에서 어떻게 구조화되는지 설명합니다.")
    add_table(s, LEFT, 1.65, [1.25, 2.15, 3.0], 0.45, ["입력", "주요 추출 항목", "사용 목적"], [
        ["결과서", "PIPE CONFIGURATION, FLOW IN PIPES, NOZZLE", "압력, 유량, 유속, 마찰손실, 구경, C-Factor 검증"],
        ["SDF", "Node, Pipe, Nozzle, Waypoint", "배관망 토폴로지, 하류 헤드 수, 네트워크 맵 구성"],
        ["CAD/DXF", "선분, 레이어, 객체 타입", "도면-아이소매트릭 대조 모듈의 선택 입력"],
        ["정책 파일", "project_meta, design_policy", "세대 내부/유입 배관, 헤드 수별 관경 정책 보완"],
    ], "입력 데이터 구조")
    add_flow(s, LEFT, 5.15, [("업로드", "결과서\nSDF"), ("파서", "표 추출\n노드 추출"), ("룰 엔진", "법적 기준\n공식 검산"), ("UI", "테이블\n맵"), ("보고서", "PDF/Excel\n리포트")], 1.02, 0.78, 0.16)

    # 3 Common
    s = blank_slide(prs)
    add_sidebar(s, 3, "공통 검증 기준", ["HW 선언", "마찰손실 재계산", "허용오차", "PASS/FAIL"], "COMMON RULE", ["공통 검증은 최적화보다 우선합니다.", "공식 재현 실패 시 결과서 확인이 필요합니다."])
    add_header(s, 3, "공통 검증 기준", "수리계산 결과가 기본 기준을 만족하는지 먼저 확인합니다.")
    add_card(s, LEFT, 1.55, 2.95, 1.22, "COMMON.HW_DECLARED", "DESIGN INFORMATION에 'Using the Hazen-Williams Equation' 선언이 있는지 확인합니다.", BLUE, SKY)
    add_card(s, LEFT + 3.25, 1.55, 2.95, 1.22, "COMMON.HW_RECALC", "Pipe별 total_length, Act. Bore, Flowrate, C-Factor로 Frict. Loss를 재계산해 비교합니다.", GREEN, GREEN_BG)
    add_table(s, LEFT, 3.65, [1.7, 2.35, 2.15], 0.42, ["항목", "수식/기준", "해석"], [
        ["총 등가길이", "length + fitting_eq + special_eq", "배관 본길이와 부속/특수설비 등가길이 합산"],
        ["HW 손실", "6.174e4·Q^1.85·L/(C^1.85·D^4.87)/0.1", "kg/cm² 기준 결과서 Frict. Loss와 비교"],
        ["허용오차", "max(0.005, reported × 0.005)", "절대오차 0.005 또는 상대 0.5% 중 큰 값"],
    ], "Hazen-Williams 재계산 기준")

    # 4 Legal
    s = blank_slide(prs)
    add_sidebar(s, 4, "법적 기준 충족 검증", ["헤드 압력/유량", "Topology 기반 유속", "가지 6 m/s", "그 밖 10 m/s"], "LEGAL CHECK", ["구경만으로 가지배관을 판단하지 않습니다.", "하류 헤드 수와 교차분기를 함께 봅니다."])
    add_header(s, 4, "법적 기준 충족 검증", "헤드와 배관 유속의 핵심 판정 기준을 정리합니다.")
    add_card(s, LEFT, 1.55, 2.95, 1.25, "헤드 검증", "표준헤드는 최소 방수압 0.1 MPa 이상, 요구 유량 이상이어야 합니다.", BLUE, SKY)
    add_card(s, LEFT + 3.25, 1.55, 2.95, 1.25, "배관 유속 검증", "PIPE/NOZZLE 구성으로 그래프를 만들고 하류 교차분기 여부로 role을 판정합니다.", ORANGE, ORANGE_BG)
    add_flow(s, LEFT, 3.55, [("Graph", "Pipe\ninput/output"), ("Nozzle", "하류 헤드\n계산"), ("Split", "교차분기\n탐지"), ("Role", "branch\nother"), ("Limit", "6 또는\n10 m/s")], 1.02, 0.78, 0.16)
    add_table(s, LEFT, 5.2, [1.25, 1.75, 3.2], 0.38, ["분류", "기준", "설명"], [
        ["branch", "6.0 m/s 이하", "50A 이하이면서 하류 subtree에 multi-head cross split이 없는 배관"],
        ["other", "10.0 m/s 이하", "50A 초과 또는 하류 subtree에 교차분기가 존재하는 배관"],
        ["review", "판정 보류", "그래프 누락, 사이클, 노드 참조 오류 등"],
    ], "Topology 기반 유속 기준")

    # 5 Pipe rules
    s = blank_slide(prs)
    add_sidebar(s, 5, "배관 항목 검증", ["PIPE.001~006", "KSD3562", "C-Factor", "REVIEW"], "PIPE RULES", ["CAD나 메타데이터가 없으면 REVIEW입니다.", "결과서는 계산 결과의 권위 출력입니다."])
    add_header(s, 5, "배관 항목 검증", "배관 재질, C-Factor, 고압 구간, 정책성 기준을 구조화합니다.")
    add_table(s, LEFT, 1.45, [1.0, 2.1, 3.15], 0.39, ["Rule", "검증 항목", "판정 방식"], [
        ["PIPE.001", "도면과 Schematic 일치", "CAD 입력이 없으면 REVIEW, 있으면 mismatch 목록 기반 판정"],
        ["PIPE.002", "1.2 MPa 이상 KSD3562", "max(inlet,outlet) ≥ 12 kg/cm² 구간은 KSD3562 요구"],
        ["PIPE.003", "재질별 C-Factor", "KSD 계열 C=120, CPVC 계열 C=150"],
        ["PIPE.004", "단위세대 내부 CPVC", "project_meta 또는 CAD zone 입력 필요"],
        ["PIPE.005", "단위세대 유입 65A 이하", "unit_inlet_pipe_labels 필요"],
        ["PIPE.006", "헤드 수별 관경 정책", "design_policy.json 정책표 기반 판정"],
    ], "4. 배관 대제목 규칙")
    add_card(s, LEFT, 6.05, 2.9, 0.72, "PASS 예시", "고압 구간이 없고 C-Factor가 재질 기준과 일치하면 배관 기본 검증은 PASS입니다.", GREEN, GREEN_BG)
    add_card(s, LEFT + 3.25, 6.05, 2.9, 0.72, "REVIEW 예시", "CAD, project_meta, design_policy가 없으면 일부 정책성 판정은 REVIEW입니다.", ORANGE, ORANGE_BG)

    # 6 Result table
    s = blank_slide(prs)
    add_sidebar(s, 6, "결과 데이터 테이블 해석", ["색상 범례", "상세 카드", "HW 검산", "Topology"], "RESULT TABLE", ["필터링된 행을 클릭하면", "수식·기준·실제값·결론이 표시됩니다."])
    add_header(s, 6, "결과 데이터 테이블 해석", "색상과 상세 카드가 어떤 의미인지 설명합니다.")
    add_image(s, LEFT, 1.3, 6.2, 2.55, "server_table.png", "결과 데이터 테이블 + 색상 범례 화면")
    add_table(s, LEFT, 4.45, [1.05, 1.6, 3.6], 0.37, ["색상", "의미", "대표 조건"], [
        ["빨강", "기준 위반", "유속 초과, HW 재계산 실패, 법적 기준 미달"],
        ["파랑", "공학 최적화 후보", "마찰손실, 변화율, 유속여유, 피팅집중, 압력여유"],
        ["초록", "경제성 검토 후보", "저유속 과설계, 압력여유 과다, 대구경 밸브"],
        ["흰색", "일반/REVIEW", "REVIEW는 별도 정보 버튼에서 규칙 상태 확인"],
    ], "색상 의미")

    # 7 Optimization
    s = blank_slide(prs)
    add_sidebar(s, 7, "설계 최적화 가이드", ["공학 최적화", "경제성 최적화", "충돌 구간", "절충 설계"], "OPTIMIZATION", ["최적화는 기준 PASS 이후 단계입니다.", "하나의 정답이 아니라 후보와 근거를 제공합니다."])
    add_header(s, 7, "설계 최적화 가이드", "공학적 안정성과 경제성 절감을 분리해서 제시합니다.")
    add_image(s, LEFT, 1.25, 6.2, 3.05, "server_optimization.png", "공학적 마찰손실 최적화 조치와 시공사 경제성 확보 방안 화면")
    add_table(s, LEFT, 4.9, [1.1, 2.55, 2.55], 0.38, ["관점", "대표 지표", "권장 조치"], [
        ["공학", "m당 손실, 변화율, 유속여유 부족", "구경 상향, 피팅 축소, 경로 단순화"],
        ["경제", "저유속, 압력여유 과다, 대구경 비용", "관경 축소 시뮬레이션, 밸브 구경 최적화"],
        ["절충", "공학 후보와 경제 후보가 충돌", "위험구간은 공학 우선, 여유 구간만 경제안 검토"],
    ], "공학/경제 관점 비교")

    # 8 Network
    s = blank_slide(prs)
    add_sidebar(s, 8, "SDF 아이소매트릭 배관망", ["줌/드래그", "적합/부적합", "Spike Map", "Economy Map"], "NETWORK VIEW", ["배관망은 판정 결과를 공간적으로 이해하기 위한 보조 화면입니다."])
    add_header(s, 8, "SDF 아이소매트릭 배관망", "검증 결과를 네트워크 맵과 연결해 해석합니다.")
    add_image(s, LEFT, 1.25, 6.2, 3.05, "server_validation.png", "검증 결과와 SDF 배관망 화면")
    add_table(s, LEFT, 4.95, [2.1, 2.0, 2.1], 0.38, ["맵", "표시 대상", "인터랙션"], [
        ["Friction Loss Spike Map", "m당 마찰손실 > 1.0 kg/cm²/m", "빨간 배관 드래그 시 정보 표시"],
        ["Economy Candidate Map", "저유속/압력여유/대구경 후보", "줌/드래그로 후보 위치 확인"],
        ["검증 결과 레이어", "선택한 PASS/FAIL 분류", "해당 노드와 엣지만 색상 표시"],
    ], "네트워크 시각화 기능")

    # 9 CAD compare
    s = blank_slide(prs)
    add_sidebar(s, 9, "CAD-아이소매트릭 대조 모듈", ["DXF 레이어 필터", "선분 그래프화", "영역 선택", "구성요소 대조"], "CAD COMPARE", ["현재 대조 모듈은 보조 검토 도구입니다.", "정확한 Pipe 단위 대조에는 도면 정제가 필요합니다."])
    add_header(s, 9, "CAD-아이소매트릭 대조 모듈", "DXF 도면과 SDF 배관망을 비교하기 위한 보조 모듈입니다.")
    add_table(s, LEFT, 1.35, [1.55, 2.25, 2.4], 0.42, ["단계", "처리 내용", "목적"], [
        ["1. DXF 정제", "배관 후보 레이어 분리, 불필요 객체 삭제", "건축 도면의 복잡도를 낮춤"],
        ["2. 선분 추출", "LINE, LWPOLYLINE, ARC endpoint 추출", "도면 선분을 네트워크 후보로 변환"],
        ["3. 점 병합", "endpoint tolerance로 인접점 병합", "노드, 말단, 분기점 생성"],
        ["4. 그래프 생성", "노드-엣지 연결 구조화", "SDF 그래프와 비교 가능한 형태 구성"],
        ["5. 대조 분석", "헤드 수, 밸브, 부속류, 길이, 연결성 비교", "명시적 mismatch 목록 생성"],
    ], "DXF 기반 대조 흐름")
    add_flow(s, LEFT, 5.7, [("DXF", "레이어\n정제"), ("Segment", "선분\n추출"), ("Node", "점 병합\n분기점"), ("Graph", "연결망\n생성"), ("Compare", "SDF와\n대조")], 1.02, 0.78, 0.16)

    # 10 Usage
    s = blank_slide(prs)
    add_sidebar(s, 10, "사용 절차와 FAQ", ["업로드 순서", "REVIEW 처리", "다운로드", "주의사항"], "FAQ", ["검증 결과는 설계자 판단을 대체하지 않습니다.", "수정 후에는 반드시 재검증해야 합니다."])
    add_header(s, 10, "사용 절차와 FAQ", "실무자가 검증 결과를 해석할 때 확인할 항목입니다.")
    add_image(s, LEFT, 1.25, 6.2, 2.35, "server_actual_main.png", "파일 업로드 및 검증 실행 첫 화면")
    add_table(s, LEFT, 4.25, [0.8, 2.4, 3.0], 0.36, ["순서", "사용자 작업", "확인 사항"], [
        ["1", "결과서 DOCX/PDF 업로드", "결과서 표 파싱 가능 여부 확인"],
        ["2", "SDF 파일 선택 업로드", "토폴로지, 배관망 맵, 하류 헤드 수 확인"],
        ["3", "검증 실행", "PASS/FAIL/WARNING/REVIEW 메시지 확인"],
        ["4", "결과 데이터 테이블 클릭", "수식, 기준값, 실제값, 판정 사유 확인"],
        ["5", "공학/경제 최적화 검토", "후보 조치 후 재계산 필요"],
    ], "기본 사용 절차")

    prs.save(OUT_ASCII)
    prs.save(OUT)
    print(OUT_ASCII, OUT_ASCII.exists(), OUT_ASCII.stat().st_size)
    print(OUT, OUT.exists(), OUT.stat().st_size)


if __name__ == "__main__":
    build()
