# -*- coding: utf-8 -*-
"""특허도면(한백수정본V1).pptx 의 텍스트를 위치 순서대로 덤프한다 (읽기 전용 조사용)."""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

SRC = Path(__file__).resolve().parents[1] / "scripts" / "특허도면(한백수정본V1).pptx"
OUT = Path(__file__).resolve().parents[1] / "data" / "_patent_dump.txt"


def collect(shapes, depth=0, acc=None):
    if acc is None:
        acc = []
    for sh in shapes:
        try:
            is_group = sh.shape_type == 6
        except Exception:
            is_group = False
        if is_group:
            collect(sh.shapes, depth + 1, acc)
            continue
        txt = ""
        if sh.has_text_frame:
            parts = []
            for para in sh.text_frame.paragraphs:
                line = "".join(r.text for r in para.runs).strip()
                if line:
                    parts.append(line)
            txt = " ⏎ ".join(parts)
        if getattr(sh, "has_table", False):
            rows = []
            for r in sh.table.rows:
                cells = [c.text.replace("\n", " ⏎ ").strip() for c in r.cells]
                rows.append(" | ".join(cells))
            txt = "[TABLE]\n    " + "\n    ".join(rows)
        if txt:
            top = Emu(sh.top).inches if sh.top is not None else 0.0
            left = Emu(sh.left).inches if sh.left is not None else 0.0
            acc.append((round(top, 2), round(left, 2), txt))
    return acc


def main() -> int:
    prs = Presentation(str(SRC))
    lines = [f"슬라이드 {len(prs.slides)}장 · 캔버스 {prs.slide_width}x{prs.slide_height}"]
    for i, slide in enumerate(prs.slides, 1):
        lines.append("")
        lines.append("=" * 72)
        lines.append(f"SLIDE {i}")
        lines.append("=" * 72)
        items = collect(slide.shapes)
        items.sort(key=lambda t: (t[0], t[1]))  # 위→아래, 왼→오른
        for top, left, txt in items:
            lines.append(f"  [{top:>6.2f},{left:>6.2f}] {txt}")
    OUT.parent.mkdir(parents=True, exist_ok=True)
    OUT.write_text("\n".join(lines), encoding="utf-8")
    print(f"wrote {OUT} ({len(lines)} lines)")
    return 0


if __name__ == "__main__":
    sys.exit(main())
